from __future__ import annotations

import os
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = Path(os.environ.get("DEFENSE_TEMPLATE", ROOT / "templates" / "defense_template.pptx"))
OUTPUT = ROOT / "output" / "defense_template_version.pptx"


SELECTED_SLIDES = [0, 1, 2, 3, 5, 6, 7, 9, 10, 25, 26, 27, 28]

NAV_MAP = {
    "基本信息": "项目概述",
    "研究背景": "证据工具",
    "研究思路": "Skill系统",
    "研究结果": "写作返修",
    "结论讨论": "图片导出",
    "创新启发": "总结展望",
    "结论与讨论": "图片与导出",
    "创新与启发": "总结与展望",
}


def remove_slide(prs: Presentation, idx: int) -> None:
    slides = prs.slides
    r_id = slides._sldIdLst[idx].rId
    prs.part.drop_rel(r_id)
    del slides._sldIdLst[idx]


def remove_shape(shape) -> None:
    shape.element.getparent().remove(shape.element)


def set_text(shape, text: str, size: float = 18, bold: bool = False, color: str = "1F2937") -> None:
    if not hasattr(shape, "text_frame"):
        return
    tf = shape.text_frame
    tf.clear()
    paragraphs = text.split("\n")
    for i, line in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor.from_string(color)
        p.space_after = Pt(2)


def add_text_box(slide, text: str, x: float, y: float, w: float, h: float, size: float = 16, bold: bool = False, color: str = "1F2937") -> None:
    shape = slide.shapes.add_textbox(Pt(x * 72), Pt(y * 72), Pt(w * 72), Pt(h * 72))
    set_text(shape, text, size=size, bold=bold, color=color)


def replace_text(shape, replacements: dict[str, str]) -> None:
    if not hasattr(shape, "text_frame"):
        return
    txt = shape.text
    if not txt:
        return
    new = txt
    for old, value in replacements.items():
        new = new.replace(old, value)
    if new != txt:
        size = 13 if len(new) < 8 else 12
        set_text(shape, new, size=size, bold=len(new) <= 8, color="1F2937")


def text_shapes(slide, *, min_x: float = 0.0, min_y: float = 0.0):
    shapes = []
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text.strip():
            x = shape.left / 914400
            y = shape.top / 914400
            if x >= min_x and y >= min_y:
                shapes.append(shape)
    return sorted(shapes, key=lambda s: (s.top, s.left))


def replace_contains(slide, needle: str, text: str, size: float = 18, bold: bool = False, color: str = "1F2937") -> None:
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and needle in shape.text:
            set_text(shape, text, size=size, bold=bold, color=color)


def clean_irrelevant_media(slide) -> None:
    for shape in list(slide.shapes):
        if shape.shape_type in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.CHART}:
            remove_shape(shape)


def remove_large_content_groups(slide) -> None:
    for shape in list(slide.shapes):
        x = shape.left / 914400
        w = shape.width / 914400
        h = shape.height / 914400
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP and x > 1.8 and w > 2.0 and h > 1.0:
            remove_shape(shape)


def apply_global_nav(prs: Presentation) -> None:
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_text(shape, NAV_MAP)


def fill_cover(slide) -> None:
    replacements = {
        "Developmental hematopoietic stem cell variation explains clonal hematopoiesis later in life":
            ("Literature Review Harness", 29, True, "111827"),
        "发育性造血干细胞变异解释了生命后期的克隆造血":
            ("面向学术综述生成的 Evidence-Grounded Agent Harness", 18, True, "111827"),
        "课题小组：组会文献汇报": ("赛题答辩：学术论文综述生成 Harness", 17, True, "111827"),
        "汇报人：XXX": ("汇报人：SurveyHarness Team", 11, False, "FFFFFF"),
        "汇报日期：20XX.XX.XX": ("汇报日期：2026.07.05", 11, False, "FFFFFF"),
    }
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            txt = shape.text.strip()
            if txt in replacements:
                text, size, bold, color = replacements[txt]
                set_text(shape, text, size=size, bold=bold, color=color)
                if text == "Literature Review Harness":
                    shape.top = Pt(1.65 * 72)
                    shape.height = Pt(0.65 * 72)
                elif text.startswith("面向学术综述"):
                    shape.top = Pt(2.75 * 72)
                    shape.height = Pt(0.45 * 72)
                elif text.startswith("赛题答辩"):
                    shape.top = Pt(3.38 * 72)
                    shape.height = Pt(0.42 * 72)


def fill_contents(slide) -> None:
    replacements = {
        "基本信息": "项目概述",
        "研究背景": "证据工具",
        "研究思路": "Skill系统",
        "研究结果": "写作返修",
        "结论与讨论": "图片与导出",
        "创新与启发": "总结与展望",
    }
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text.strip() in replacements:
            set_text(shape, replacements[shape.text.strip()], size=18, bold=True, color="1F2937")


def fill_project_overview(slide) -> None:
    replace_contains(slide, "1 项目概述", "1 项目概述", size=20, bold=True, color="1F2937")
    for shape in list(slide.shapes):
        x = shape.left / 914400
        y = shape.top / 914400
        if 2.0 <= x <= 5.7 and 3.0 <= y <= 4.45 and not shape.has_table:
            remove_shape(shape)
    table_values = [
        ("项目", "Literature Review Harness"),
        ("任务", "学术论文综述生成"),
        ("主题", "World Models Demo"),
        ("核心", "Evidence-grounded Agent"),
        ("输出", "MD / HTML / LaTeX / Figures"),
    ]
    for shape in slide.shapes:
        if shape.has_table:
            for row, (k, v) in zip(shape.table.rows, table_values):
                row.cells[0].text = k
                row.cells[1].text = v
                for j, cell in enumerate(row.cells):
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = "Microsoft YaHei"
                            run.font.size = Pt(11)
                            run.font.bold = j == 0
                            run.font.color.rgb = RGBColor.from_string("1F2937")
    replace_contains(slide, "主要 | 研究内容", "核心 | 项目定位", size=17, bold=True, color="FFFFFF")
    replace_contains(
        slide,
        "请输入文献主要研究内容",
        "本项目实现一个可复用的学术综述生成 Harness：先用 Sciverse 检索论文片段，必要时用 MinerU 解析全文，再构建 Evidence KB；随后由 AgentLoop 调用 skill、生成 survey context、完成写作、审查、引用校验、图片生成和多格式导出。",
        size=15.5,
        bold=False,
        color="1F2937",
    )


def fill_background(slide) -> None:
    remove_large_content_groups(slide)
    replace_contains(slide, "2.1 证据工具", "2.1 赛题背景与核心问题", size=20, bold=True)
    add_text_box(
        slide,
        "赛题要求\n\n构建论文综述 Harness，能够阅读论文、分类整理、总结关键论文、梳理发展脉络、提出未来研究方向，并生成图文并茂综述。\n\n核心底线：无幻觉引用，无缺乏事实证据的描述。",
        2.25, 1.55, 4.65, 4.8, size=14.6, color="1F2937",
    )
    replace_contains(slide, "造血干细胞", "为什么不能直接写？", size=17, bold=True, color="1F2937")
    replace_contains(
        slide,
        "造血干细胞（",
        "直接让 LLM 写综述，容易出现编造论文、引用不可追踪、结构像摘要堆叠、返修不可控等问题。赛题强调无幻觉引用，因此需要把每个论断尽量绑定到真实证据。",
        size=15.5,
    )
    replace_contains(slide, "克隆性造血", "Harness 的设计目标", size=17, bold=True, color="1F2937")
    replace_contains(
        slide,
        "克隆性造血是",
        "系统不追求一次性生成，而是把检索、证据、规划、写作、审查、校验和导出拆成可审计流程。World Models 是展示题目，框架可以迁移到其他领域。",
        size=15.5,
    )


def fill_logic(slide) -> None:
    replace_contains(slide, "3 Skill系统", "3 研究思路：Agent Harness 主链路", size=20, bold=True)
    contents = [
        "用户输入综述主题：World Models / 其他学术领域",
        "Sciverse 检索论文片段、URL、元数据和相关性分数",
        "MinerU 按需解析关键论文全文，用于补读细节",
        "Evidence KB 统一 paper、evidence、parsed document",
        "Skill 系统按阶段发现、路由、加载、卸载写作协议",
        "Survey Context 生成 outline、timeline、citation map",
        "AgentLoop 输出草稿，并进入 reviewer + verifier 返修",
        "导出 Markdown / HTML / LaTeX / 图片和审计文件",
    ]
    shapes = [s for s in text_shapes(slide, min_x=2.0) if s.text.strip() and "3 Skill系统" not in s.text]
    for shape, content in zip(shapes, contents):
        set_text(shape, content, size=12.5, bold=False, color="1F2937")


def fill_evidence(slide) -> None:
    remove_large_content_groups(slide)
    replace_contains(slide, "4 写作返修1", "4.1 证据层：Sciverse + MinerU + LiteratureKB", size=20, bold=True)
    add_text_box(
        slide,
        "证据链路\n\n1. Sciverse 召回候选论文和片段。\n2. Evidence KB 记录 paper / evidence / source。\n3. 关键论文再由 MinerU 补全文解析。\n4. 正文只引用 KB 中存在的 evidence id。",
        2.35, 1.25, 4.75, 5.2, size=14.7, color="1F2937",
    )
    replace_contains(slide, "多种特殊变异", "Sciverse 负责广度检索：返回论文片段、元数据、URL、score，快速建立候选证据池。", size=14.2)
    replace_contains(slide, "借助CRISPR", "MinerU 负责深度补读：关键论文需要更多上下文时，Agent 可读取 parsed paper 原文解析结果。LiteratureKB 生成稳定 paper_id 和 evidence_id，正文引用只允许使用真实 evidence id。", size=14.2)


def fill_tools(slide) -> None:
    remove_large_content_groups(slide)
    replace_contains(slide, "4 写作返修1", "4.2 工具层：模型在边界内自主调用", size=20, bold=True)
    add_text_box(
        slide,
        "AgentLoop 职责\n\n- 构造模型上下文\n- 暴露工具 schema\n- 执行 tool calls\n- 写回工具结果\n- 控制迭代和返修预算\n- 保存运行审计产物",
        2.35, 1.35, 4.75, 5.0, size=14.7, color="1F2937",
    )
    replace_contains(
        slide,
        "对WT野生型组",
        "核心工具：build_literature_kb / search_literature / list_evidence / read_evidence / list_parsed_papers / read_parsed_paper / prepare_survey_context。",
        size=14.8,
    )
    replace_contains(
        slide,
        "把WT野生型",
        "模型通过工具读材料、补证据、规划结构，再输出综述；Harness 负责状态、权限、trace、失败降级和返修预算。",
        size=14.8,
    )


def fill_skill(slide) -> None:
    remove_large_content_groups(slide)
    replace_contains(slide, "4 写作返修2", "4.3 Skill System：渐进式披露", size=20, bold=True)
    add_text_box(
        slide,
        "Skill 调用逻辑\n\n发现：只读 metadata，避免上下文爆炸。\n路由：按 phase / role 选择少量 skill。\n加载：只注入当前阶段需要的 SKILL.md。\n资源：必要时读取 skill 附带文件或脚本。\n卸载：阶段结束后清理 active skill。",
        2.35, 1.35, 6.2, 4.9, size=14.5, color="1F2937",
    )
    replace_contains(
        slide,
        "研究结果1明确指出",
        "SkillManager 统一扫描 skills/；list_index 只暴露短描述、phase、roles；route_for_phase 选择当前阶段需要的 skill；load_for_phase 注入 SKILL.md；unload 清理上下文，并写入 skill_trace.json。",
        size=14.2,
    )
    replace_contains(
        slide,
        "为探究PR-DUB",
        "设计目的：让外部写作/引用/导出 skill 可以被 Agent 发现和使用，但不把所有 skill 内容一次塞满上下文窗口。",
        size=14.2,
    )


def fill_review(slide) -> None:
    remove_large_content_groups(slide)
    replace_contains(slide, "4 写作返修2", "4.4 写作、审查与引用校验", size=20, bold=True)
    add_text_box(
        slide,
        "三路并行 Reviewer\n\n内容质量：是否像 survey，而不是摘要堆叠。\n引用准确：是否每个关键论断都有证据。\n结构完整：是否包含背景、谱系、比较、局限和研究议程。\n\n任一失败：反馈写回消息历史，进入下一轮修复。",
        2.35, 1.35, 6.1, 4.95, size=14.2, color="1F2937",
    )
    replace_contains(
        slide,
        "通过基因富集分析",
        "LLM 输出完整 survey 后，三路 reviewer 并行审查：内容质量、引用准确性、结构完整性。任一不通过，就把 issues 和 suggestions 注入下一轮消息，触发修复。",
        size=14.2,
    )
    replace_contains(
        slide,
        "相关修复蛋白",
        "CitationVerifier 进行机械校验：未知 evidence id、长段落缺 citation、工具日志污染、References 不一致都会进入 check_report 并触发返修。",
        size=14.2,
    )


def fill_outputs(slide) -> None:
    remove_large_content_groups(slide)
    replace_contains(slide, "5.1 结论", "5.1 图片生成与输出包", size=20, bold=True)
    replace_contains(slide, "5.1 图片导出", "5.1 图片生成与输出包", size=20, bold=True)
    add_text_box(
        slide,
        "输出产物\n\nsurvey.md / survey.html / survey.tex\n\nevidence_pack.json\ncheck_report.json\nskill_trace.json\nfigure_plan.json\nfigures/",
        2.35, 1.35, 4.7, 5.35, size=14.5, color="1F2937",
    )
    replace_contains(
        slide,
        "综上所述",
        "图片生成不是独立装饰，而是根据 survey 章节和 evidence ids 规划。Figure planner 控制数量，OpenAI-compatible image API 生成关键图；时间线、矩阵等结构图可用 SVG fallback。\n\n每次运行输出到 output/runs/YYYYMMDD-HHMM-topic/，包括 survey.md、survey.html、survey.tex、evidence_pack.json、check_report.json、skill_trace.json、figure_plan.json 和 figures/。",
        size=14.2,
    )


def fill_limitations(slide) -> None:
    replace_contains(slide, "5.2 不足之处", "5.2 当前局限与改进方向", size=20, bold=True)
    items = [
        ("长文质量", "仍受模型能力影响，需要更强的写作 prompt 和 reviewer patch plan。"),
        ("检索质量", "依赖 Sciverse 返回结果，需要更强 ranking、dedup 和 bibliography formatter。"),
        ("全文解析", "MinerU 依赖 URL 与解析时间预算，需要更细的补读策略。"),
        ("论文格式", "LaTeX 目前是基础导出，后续接入 arXiv / conference 模板。"),
    ]
    headers = [s for s in text_shapes(slide, min_x=2.0) if "不足之处" in s.text]
    bodies = [s for s in text_shapes(slide, min_x=2.0) if "请输入正文内容" in s.text]
    for shape, (head, _) in zip(headers, items):
        set_text(shape, head, size=15.5, bold=True, color="1F2937")
    for shape, (_, body) in zip(bodies, items):
        set_text(shape, body, size=12.5, bold=False, color="1F2937")


def fill_innovation(slide) -> None:
    replace_contains(slide, "6 总结展望", "6 创新与启发", size=20, bold=True)
    replacements = [
        ("Evidence-first：先构建证据，再写正文，减少幻觉引用。", "Agentic：模型通过工具完成检索、阅读、规划和返修。", "Verified：三路 reviewer + CitationVerifier 形成质量门禁。"),
        ("Skill-aware：外部 skill 可发现、可路由、可加载、可卸载。", "Traceable：evidence_pack、check_report、skill_trace 可复查。", "Extensible：后续可接入 memory、MCP、多 agent 与正式 LaTeX 模板。"),
    ]
    text_blocks = [s for s in text_shapes(slide, min_x=3.0) if "请输入正文内容" in s.text]
    flat = [x for group in replacements for x in group]
    for shape, text in zip(text_blocks, flat):
        set_text(shape, text, size=12.2, bold=False, color="1F2937")


def fill_closing(slide) -> None:
    replace_contains(slide, "汇报完毕", "汇报完毕，敬请老师同学批评指正！", size=28, bold=True, color="FFFFFF")
    replace_contains(slide, "20XX年XX大学XX学院专业", "Literature Review Harness · Evidence-grounded Survey Generation", size=14, bold=False, color="FFFFFF")
    replace_contains(slide, "汇报人：XXX", "汇报人：SurveyHarness Team", size=13, bold=False, color="1F2937")
    replace_contains(slide, "汇报日期：20XX.XX.XX", "汇报日期：2026.07.05", size=13, bold=False, color="1F2937")


def main() -> None:
    prs = Presentation(TEMPLATE)
    for idx in reversed(range(len(prs.slides))):
        if idx not in SELECTED_SLIDES:
            remove_slide(prs, idx)

    for slide in prs.slides:
        clean_irrelevant_media(slide)

    apply_global_nav(prs)

    fill_cover(prs.slides[0])
    fill_contents(prs.slides[1])
    fill_project_overview(prs.slides[2])
    fill_background(prs.slides[3])
    fill_logic(prs.slides[4])
    fill_evidence(prs.slides[5])
    fill_tools(prs.slides[6])
    fill_skill(prs.slides[7])
    fill_review(prs.slides[8])
    fill_outputs(prs.slides[9])
    fill_limitations(prs.slides[10])
    fill_innovation(prs.slides[11])
    fill_closing(prs.slides[12])

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
